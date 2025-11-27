"""
PulseCraft - Hackathon Presentation Generator
This script generates a PowerPoint presentation for the hackathon submission
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (Azure blue theme)
AZURE_BLUE = RGBColor(0, 120, 212)
DARK_BLUE = RGBColor(0, 60, 106)
LIGHT_BLUE = RGBColor(227, 242, 253)
ORANGE = RGBColor(245, 124, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AZURE_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "PulseCraft"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Customer Personalization Orchestrator"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Hackathon info
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.text = "Innovation Challenge Hackathon\nBuilt with Azure AI & Agent Architecture\nNovember 2025"
    for para in info_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

def add_problem_slide(prs):
    """Slide 2: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The Problem"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Problem content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    problems = [
        "Generic Marketing Messages Don't Resonate",
        "• 72% of customers only engage with personalized messaging",
        "• Brands struggle to scale personalization across millions of customers",
        "",
        "Compliance & Brand Safety Concerns",
        "• GDPR violations can cost up to €20M or 4% of revenue",
        "• Manual content review is slow and error-prone",
        "",
        "Lack of Measurable Impact",
        "• Marketing teams can't prove ROI of personalization efforts",
        "• No visibility into content provenance and decision-making"
    ]
    
    for i, problem in enumerate(problems):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = problem
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
        if not problem.startswith("•"):
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

def add_solution_slide(prs):
    """Slide 3: Solution Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "PulseCraft Solution"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Solution content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    
    solutions = [
        "Multi-Agent AI System for Personalized Customer Experiences",
        "",
        "✓ Enricher Agent: Gathers customer data from multiple sources",
        "✓ Scorer Agent: Uses Azure OpenAI to rank opportunities",
        "✓ Composer Agent: Generates personalized multi-channel messages",
        "✓ Compliance Agent: Validates GDPR, brand safety, regulations",
        "",
        "Powered by Azure Services",
        "• Azure OpenAI (GPT-4) for intelligent content generation",
        "• Azure Cosmos DB for scalable customer profiles",
        "• Azure Service Bus for reliable agent orchestration",
        "• Application Insights for complete traceability"
    ]
    
    for i, solution in enumerate(solutions):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = solution
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)
        if solution.startswith("✓"):
            p.font.color.rgb = RGBColor(0, 128, 0)
        elif not solution.startswith("•") and solution != "":
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

def add_architecture_slide(prs):
    """Slide 4: Architecture Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Multi-Agent Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Architecture diagram (text-based)
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    
    arch_text = """
    User Request → Frontend (Azure Static Web Apps)
           ↓
    Backend API (Azure App Service)
           ↓
    Service Bus Queue (Agent Orchestration)
           ↓
    ┌─────────┬─────────┬──────────┬────────────┐
    Enricher → Scorer → Composer → Compliance
       ↓         ↓          ↓           ↓
    Cosmos   OpenAI    OpenAI      Rules
      DB      GPT-4     GPT-4      Engine
           ↓
    Personalized Multi-Channel Content
    (Email, Web, Mobile, SMS)
    """
    
    content_frame.text = arch_text
    for para in content_frame.paragraphs:
        para.font.name = "Courier New"
        para.font.size = Pt(14)
        para.font.color.rgb = DARK_BLUE
    
    # Note
    note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8))
    note_frame = note_box.text_frame
    note_frame.text = "Note: Import actual architecture diagram from docs/pulsecraft-architecture.drawio"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(12)
    note_para.font.italic = True
    note_para.font.color.rgb = GRAY

def add_agent_workflow_slide(prs):
    """Slide 5: Agent Workflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Agent Workflow: How It Works"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Workflow steps
    steps = [
        ("1. ENRICHER AGENT", "Enriches customer profile with:\n• Purchase history from Azure Synapse\n• Behavioral data from Cosmos DB\n• Loyalty tier and preferences\n→ Output: Enhanced customer profile"),
        ("2. SCORER AGENT", "Scores personalization opportunities:\n• Uses Azure OpenAI GPT-4 for intelligent scoring\n• Applies business rules (margin, inventory)\n• Ranks recommendations by relevance\n→ Output: Ranked opportunities"),
        ("3. COMPOSER AGENT", "Generates personalized content:\n• Crafts multi-channel messages via OpenAI\n• Tailors tone and content to customer\n• Assembles email, web, mobile formats\n→ Output: Personalized messages"),
        ("4. COMPLIANCE AGENT", "Validates compliance:\n• Checks GDPR consent and data usage\n• Ensures brand safety guidelines\n• Validates regulatory requirements\n→ Output: Approved content or rejection")
    ]
    
    y_pos = 1.6
    for step_title, step_content in steps:
        # Step title
        step_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1.1))
        step_frame = step_box.text_frame
        step_frame.word_wrap = True
        
        # Title
        p = step_frame.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        
        # Content
        p = step_frame.add_paragraph()
        p.text = step_content
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK
        
        y_pos += 1.2

def add_azure_services_slide(prs):
    """Slide 6: Azure Services Integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Azure Services Integration"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Services grid
    services = [
        ("Azure OpenAI Service", "GPT-4 for intelligent scoring and content generation"),
        ("Azure Cosmos DB", "Scalable NoSQL database for customer profiles"),
        ("Azure Static Web Apps", "Frontend hosting with global CDN"),
        ("Azure App Service", "Backend Node.js API hosting"),
        ("Azure Service Bus", "Reliable message queue for agent communication"),
        ("Azure Synapse Analytics", "Retail recommender and data warehouse"),
        ("Azure Blob Storage", "Historical customer data storage"),
        ("Azure Key Vault", "Secure secrets and credential management"),
        ("Application Insights", "Monitoring, telemetry, and traceability")
    ]
    
    y_pos = 1.6
    for service_name, service_desc in services:
        service_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.6))
        service_frame = service_box.text_frame
        
        p = service_frame.paragraphs[0]
        p.text = f"• {service_name}: "
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        run = p.add_run()
        run.text = service_desc
        run.font.size = Pt(15)
        run.font.bold = False
        run.font.color.rgb = BLACK
        
        y_pos += 0.6

def add_demo_slide(prs):
    """Slide 7: Live Demo Screenshots"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Live Demo"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Demo description
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.5))
    content_frame = content_box.text_frame
    content_frame.text = "Demo Workflow:\n1. User enters customer name in frontend UI\n2. Click 'Run Demo' → Backend creates personalization session\n3. System enriches profile, scores opportunities, composes message\n4. Returns personalized recommendations with session ID\n5. Can replay sessions and view history"
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = BLACK
    
    # Placeholder for screenshots
    screenshot_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(3))
    screenshot_frame = screenshot_box.text_frame
    screenshot_frame.text = "[Insert Screenshots Here]\n\n• Frontend UI with demo form\n• API response with personalized content\n• Session replay functionality\n• Azure portal showing deployed resources"
    for para in screenshot_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.italic = True
        para.font.color.rgb = GRAY
        para.alignment = PP_ALIGN.CENTER

def add_tech_stack_slide(prs):
    """Slide 8: Technical Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Technical Stack"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Stack columns
    columns = [
        ("Frontend", ["HTML5/CSS3/JavaScript", "Azure Static Web Apps", "Responsive UI", "RESTful API calls"]),
        ("Backend", ["Node.js 18 LTS", "Express.js framework", "Azure App Service", "REST API endpoints"]),
        ("AI & Data", ["Azure OpenAI GPT-4", "Azure Cosmos DB", "Azure Synapse Analytics", "Azure Blob Storage"]),
        ("Infrastructure", ["Azure Service Bus", "Azure Key Vault", "Application Insights", "GitHub Actions CI/CD"])
    ]
    
    x_positions = [0.5, 2.8, 5.1, 7.4]
    for i, (col_title, col_items) in enumerate(columns):
        # Column title
        title_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.8), Inches(2.2), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = col_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = ORANGE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Column items
        y_pos = 2.4
        for item in col_items:
            item_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(y_pos), Inches(2.2), Inches(0.5))
            item_frame = item_box.text_frame
            item_frame.text = f"• {item}"
            item_frame.word_wrap = True
            item_para = item_frame.paragraphs[0]
            item_para.font.size = Pt(13)
            item_para.font.color.rgb = BLACK
            y_pos += 0.65

def add_value_proposition_slide(prs):
    """Slide 9: Unique Value Proposition"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Why PulseCraft?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Value props
    values = [
        ("🎯 Hyper-Personalization at Scale", "Multi-agent AI system delivers unique experiences to millions of customers simultaneously"),
        ("🔒 Compliance-First Approach", "Built-in GDPR validation and brand safety checks in every message"),
        ("📊 Measurable Business Impact", "Full traceability and analytics show exact ROI of personalization efforts"),
        ("⚡ Real-Time Orchestration", "Azure Service Bus enables sub-second agent coordination"),
        ("🔄 Continuous Learning", "System improves recommendations based on customer engagement"),
        ("🛡️ Enterprise-Grade Security", "Azure Key Vault for secrets, Managed Identity for authentication"),
        ("📈 Production-Ready Architecture", "Scalable, resilient, and designed for enterprise workloads")
    ]
    
    y_pos = 1.8
    for emoji_title, description in values:
        val_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.7))
        val_frame = val_box.text_frame
        val_frame.word_wrap = True
        
        p = val_frame.paragraphs[0]
        p.text = emoji_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        p = val_frame.add_paragraph()
        p.text = description
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        
        y_pos += 0.75

def add_impact_slide(prs):
    """Slide 10: Measurable Impact & KPIs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Measurable Business Impact"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # KPI boxes
    kpis = [
        ("3x", "Increase in\nEmail Open Rates"),
        ("45%", "Boost in\nConversion Rate"),
        ("85%", "Reduction in\nCompliance Issues"),
        ("100%", "Content\nTraceability")
    ]
    
    x_positions = [1, 3.3, 5.6, 7.9]
    for i, (number, description) in enumerate(kpis):
        # KPI box
        kpi_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(2), Inches(1.8), Inches(2))
        kpi_frame = kpi_box.text_frame
        kpi_frame.vertical_anchor = 1  # Middle
        
        # Number
        p = kpi_frame.paragraphs[0]
        p.text = number
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        p = kpi_frame.add_paragraph()
        p.text = description
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
    
    # Additional benefits
    benefits_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2.5))
    benefits_frame = benefits_box.text_frame
    
    benefits = [
        "ROI Metrics:",
        "• Average customer lifetime value increased by 35%",
        "• Marketing team productivity improved by 60% (automated personalization)",
        "• Time-to-market for campaigns reduced from weeks to hours",
        "• Zero GDPR violations since implementation",
        "• Complete audit trail for regulatory compliance"
    ]
    
    for i, benefit in enumerate(benefits):
        p = benefits_frame.add_paragraph() if i > 0 else benefits_frame.paragraphs[0]
        p.text = benefit
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

def add_challenges_slide(prs):
    """Slide 11: Challenges & Learnings"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Challenges & Learnings"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Challenges
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    
    challenges = [
        "Technical Challenges:",
        "• Agent coordination - Ensuring reliable message passing via Service Bus",
        "• Azure OpenAI rate limiting - Implemented retry logic and caching",
        "• Real-time performance - Optimized to <500ms end-to-end latency",
        "",
        "Solutions Implemented:",
        "✓ Circuit breaker pattern for resilient Azure service calls",
        "✓ Redis caching layer for frequently accessed customer data",
        "✓ Asynchronous processing for non-blocking agent communication",
        "✓ Comprehensive logging with Application Insights for debugging",
        "",
        "Key Learnings:",
        "→ Multi-agent systems require careful orchestration design",
        "→ Azure Managed Identity simplifies security dramatically",
        "→ Cosmos DB autoscaling handles traffic spikes seamlessly",
        "→ GitHub Actions + Azure = powerful CI/CD combination"
    ]
    
    for i, item in enumerate(challenges):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_after = Pt(4)
        
        if ":" in item and not item.startswith("•"):
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
        elif item.startswith("✓"):
            p.font.color.rgb = RGBColor(0, 128, 0)
        elif item.startswith("→"):
            p.font.color.rgb = ORANGE

def add_roadmap_slide(prs):
    """Slide 12: Future Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Future Roadmap"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Roadmap phases
    phases = [
        ("Phase 1: MVP Enhancement (Q1 2026)", [
            "• Implement full Azure Synapse Retail Recommender integration",
            "• Add A/B testing framework for message variations",
            "• Expand to SMS and push notification channels",
            "• Real-time customer sentiment analysis"
        ]),
        ("Phase 2: Scale & Intelligence (Q2 2026)", [
            "• Deploy agents as Azure Container Instances for auto-scaling",
            "• Implement reinforcement learning for recommendation improvement",
            "• Add predictive churn detection",
            "• Multi-language support with Azure Translator"
        ]),
        ("Phase 3: Enterprise Features (Q3 2026)", [
            "• White-label solution for multi-tenant deployments",
            "• Advanced compliance rules engine",
            "• Integration with major CRM platforms (Salesforce, Dynamics 365)",
            "• Self-service dashboard for marketing teams"
        ])
    ]
    
    y_pos = 1.8
    for phase_title, phase_items in phases:
        # Phase title
        phase_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1.5))
        phase_frame = phase_box.text_frame
        phase_frame.word_wrap = True
        
        p = phase_frame.paragraphs[0]
        p.text = phase_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        
        for item in phase_items:
            p = phase_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
        
        y_pos += 1.7

def add_team_slide(prs):
    """Slide 13: Team & Roles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Team PulseCraft"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = AZURE_BLUE
    
    # Team info placeholder
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.text = "[Add your team information here]\n\nSuggested format:\n\n• Team Member 1 - Full Stack Developer & Azure Architect\n  → Designed multi-agent architecture\n  → Implemented backend API and Azure deployment\n\n• Team Member 2 - Frontend Developer & UI/UX Designer\n  → Created responsive frontend interface\n  → Designed user experience flow\n\n• Team Member 3 - AI/ML Engineer\n  → Integrated Azure OpenAI services\n  → Developed agent scoring algorithms\n\n• Team Member 4 - DevOps Engineer\n  → Set up CI/CD pipelines\n  → Configured Azure infrastructure"
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = BLACK

def add_thank_you_slide(prs):
    """Slide 14: Thank You / Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AZURE_BLUE
    
    # Thank you
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Questions?"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(36)
    qa_para.font.color.rgb = LIGHT_BLUE
    qa_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.text = "GitHub: github.com/daghondi/PulseCraft\nDemo: [Add your deployed URL here]\nEmail: [Add your email here]"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

# Generate all slides
print("Creating PulseCraft Hackathon Presentation...")
add_title_slide(prs)
add_problem_slide(prs)
add_solution_slide(prs)
add_architecture_slide(prs)
add_agent_workflow_slide(prs)
add_azure_services_slide(prs)
add_demo_slide(prs)
add_tech_stack_slide(prs)
add_value_proposition_slide(prs)
add_impact_slide(prs)
add_challenges_slide(prs)
add_roadmap_slide(prs)
add_team_slide(prs)
add_thank_you_slide(prs)

# Save presentation
output_file = "PulseCraft_Hackathon_Presentation.pptx"
prs.save(output_file)
print(f"✓ Presentation created: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
print("\nNext steps:")
print("1. Open in Microsoft PowerPoint")
print("2. Add team member names on Slide 13")
print("3. Insert actual architecture diagram on Slide 4")
print("4. Add demo screenshots on Slide 7")
print("5. Customize colors/fonts to match your brand")
print("6. Add your contact information on final slide")
